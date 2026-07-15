# -*- coding: utf-8 -*-
"""Genera LA presentación del Comparador Regulatorio de Calidad de Gas (estilo Enagás).

    python docs/generar_presentacion.py

Presentación de apoyo (~5 minutos, 7 diapositivas). Se expone DESPUÉS de la demo en vivo
del programa, así que no repite el "qué hace": resume, y desarrolla lo que la demo no
enseña —la ARQUITECTURA y la ONTOLOGÍA—.

  1. Portada
  2. Resumen: qué es y qué garantiza
  3. ARQUITECTURA — el esquema
  4. ARQUITECTURA — las tres capas de datos
  5. ONTOLOGÍA — la estructura real del fichero
  6. ONTOLOGÍA — anatomía de un valor y estados de verificación
  7. Cierre — las garantías

Los esquemas NO son imágenes: se dibujan con formas nativas de PowerPoint, adaptadas al
formato 16:9, para que se lean bien y sean editables.
TODAS las cifras se leen de la ontología al generar (nunca se teclean a mano).
Incluye NOTAS DEL PONENTE en cada diapositiva.

Salida: docs/Presentacion_Comparador_Gas.pptx
"""
import os
import io
import collections
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

AQUI = os.path.dirname(os.path.abspath(__file__))
RAIZ = os.path.dirname(AQUI)
ONT = os.path.join(RAIZ, "data", "ontologia", "ontologia_enagas.yaml")
PPTX = os.path.join(AQUI, "Presentacion_Comparador_Gas.pptx")

# ----------------------- Datos reales leídos de la ontología -----------------------
_d = yaml.safe_load(io.open(ONT, encoding="utf-8"))
ONTO = _d["ontologia"]
VER_ONT = ONTO.get("version", "—")
REV_ONT = ONTO.get("fecha_revision", "—")
N_FUENTES = len(ONTO["fuentes_normativas"])
N_GASES = len(ONTO.get("tipos_gas", []) or [])
CLAVES_RAIZ = list(_d.keys())


def _resumen(seccion):
    """(nº parámetros, nº jurisdicciones, nº celdas, Counter de estados) de una sección."""
    P = _d.get(seccion) or {}
    est = collections.Counter()
    jur = set()
    for v in P.values():
        for j, lim in (v.get("limites") or {}).items():
            jur.add(j)
            est[(lim or {}).get("estado_verificacion", "?")] += 1
    return len(P), len(jur), sum(est.values()), est


GN_P, GN_J, GN_C, GN_E = _resumen("parametros")
BIO_P, BIO_J, BIO_C, BIO_E = _resumen("parametros_biometano")
H2_P, H2_J, H2_C, H2_E = _resumen("parametros_hidrogeno")

V = "VERIFICADO"
VS = "VERIFICADO_SECUNDARIO"
NV = "NO_VERIFICABLE_SIN_FUENTE"
N_VERIF, N_NOVER = GN_E[V], GN_E[NV]
TOT_CELDAS = GN_C + BIO_C + H2_C
TOT_V = GN_E[V] + BIO_E[V] + H2_E[V]
TOT_VS = GN_E[VS] + BIO_E[VS] + H2_E[VS]
TOT_NV = GN_E[NV] + BIO_E[NV] + H2_E[NV]

# Claves que el código NO lee: son documentación de diseño dentro del YAML.
CLAVES_INERTES = [k for k in ("flags", "orquestador", "motor_reglas", "rag") if k in _d]

# ----------------------------- Paleta -----------------------------
AZUL = RGBColor(0x01, 0x3A, 0x57)
AZUL2 = RGBColor(0x0A, 0x5A, 0x82)
CYAN = RGBColor(0x00, 0x99, 0xD6)
VERDE = RGBColor(0x6C, 0xB3, 0x3E)
VERDEOS = RGBColor(0x3E, 0x6B, 0x22)
VERDECL = RGBColor(0xE8, 0xF4, 0xE2)
NARANJA = RGBColor(0xE8, 0x8A, 0x1A)
NARANJAOS = RGBColor(0xA5, 0x61, 0x0E)
NARANJACL = RGBColor(0xFD, 0xF0, 0xDF)
GRIS = RGBColor(0x4A, 0x5B, 0x68)
TINTA = RGBColor(0x1B, 0x2A, 0x38)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRISCL = RGBColor(0xEE, 0xF2, 0xF6)
BORDE = RGBColor(0xC8, 0xD4, 0xDD)
AZULCL = RGBColor(0xE4, 0xEF, 0xF6)
FUENTE = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _slide():
    return prs.slides.add_slide(BLANK)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _shape(slide, forma, l, t, w, h, fill=None, line=None, line_w=None, dash=False):
    shp = slide.shapes.add_shape(forma, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w or 1)
        if dash:
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shp.shadow.inherit = False
    if shp.has_text_frame:
        shp.text_frame.word_wrap = True
    return shp


def _box(slide, l, t, w, h, **kw):
    return _shape(slide, MSO_SHAPE.RECTANGLE, l, t, w, h, **kw)


def _rbox(slide, l, t, w, h, **kw):
    shp = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, **kw)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


def _flecha(slide, l, t, w, h, sentido="der", color=None):
    formas = {"der": MSO_SHAPE.RIGHT_ARROW, "izq": MSO_SHAPE.LEFT_ARROW,
              "arr": MSO_SHAPE.UP_ARROW, "aba": MSO_SHAPE.DOWN_ARROW}
    return _shape(slide, formas[sentido], l, t, w, h, fill=color or CYAN)


def _texto(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6, font=FUENTE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, r in enumerate(runs):
        txt, size, color, bold, bullet, level = (r + (None,) * 6)[:6]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.level = level or 0
        run = p.add_run(); run.text = ("• " + txt) if bullet else txt
        run.font.name = font; run.font.size = Pt(size)
        run.font.color.rgb = color; run.font.bold = bool(bold)
    return tb


def _cabecera(slide, titulo, n, kicker=None):
    _box(slide, 0, 0, 13.333, 1.02, fill=AZUL)
    _box(slide, 0, 1.02, 13.333, 0.055, fill=CYAN)
    ty = 0.16
    if kicker:
        _texto(slide, 0.55, 0.13, 11.5, 0.26, [(kicker.upper(), 11, CYAN, True, False, 0)])
        ty = 0.38
    _texto(slide, 0.55, ty, 11.6, 0.58, [(titulo, 25, BLANCO, True, False, 0)],
           anchor=MSO_ANCHOR.MIDDLE if not kicker else MSO_ANCHOR.TOP)
    _texto(slide, 0.55, 7.06, 9, 0.32,
           [("Comparador Regulatorio de Calidad de Gas — Enagás", 10, GRIS, False, False, 0)])
    _texto(slide, 12.2, 7.06, 0.85, 0.32, [(str(n), 10, GRIS, False, False, 0)], align=PP_ALIGN.RIGHT)


N = 0


def nxt():
    global N
    N += 1
    return N


# ===================== CONTENIDO — 7 diapositivas (~5 min) =====================

# ------------------------------ 1. Portada ------------------------------
s = _slide()
_box(s, 0, 0, 13.333, 7.5, fill=AZUL)
_box(s, 0, 4.62, 13.333, 0.06, fill=CYAN)
_box(s, 0.95, 2.05, 0.15, 2.25, fill=VERDE)
_texto(s, 1.35, 2.0, 11, 1.2, [
    ("Comparador Regulatorio", 42, BLANCO, True, False, 0),
    ("de Calidad de Gas", 42, BLANCO, True, False, 0),
], space=2)
_texto(s, 1.38, 4.85, 11, 1.2, [
    ("Lo que hay debajo: arquitectura y ontología", 21, CYAN, False, False, 0),
    (f"Acabáis de ver QUÉ hace. Ahora, CÓMO está construido y POR QUÉ es fiable.",
     15, RGBColor(0xC9, 0xD6, 0xE0), False, False, 0),
], space=7)
_texto(s, 1.38, 6.55, 11, 0.4,
       [(f"Ontología v{VER_ONT} · rev. {REV_ONT}", 12, RGBColor(0x8F, 0xA6, 0xB8), False, False, 0)])
_notes(s, "[~15 s] Enlazar con la demo: 'Acabáis de ver lo que hace. En cinco minutos os enseño lo que hay "
          "debajo: la arquitectura y, sobre todo, la ontología, que es donde vive el rigor del sistema.' "
          "No repitáis funcionalidades: ya las han visto.")

# ------------------- 2. Resumen: qué es y qué garantiza -------------------
s = _slide(); _cabecera(s, "El sistema, en una diapositiva", nxt(), kicker="Resumen")
_texto(s, 0.6, 1.35, 12.2, 0.35,
       [("Lo que acabáis de ver, resumido — y la promesa que lo sostiene.", 14, GRIS, False, False, 0)])
for num, txt, color, x in ((str(GN_J), "jurisdicciones", CYAN, 0.6),
                           (str(GN_P), "parámetros de calidad", CYAN, 3.71),
                           (str(TOT_CELDAS), "registros trazables", CYAN, 6.82),
                           ("0", "cifras inventadas", VERDE, 9.93)):
    _box(s, x, 1.8, 2.95, 1.75, fill=GRISCL)
    _box(s, x, 1.8, 2.95, 0.08, fill=color)
    _texto(s, x, 2.0, 2.95, 0.9, [(num, 42, AZUL, True, False, 0)], align=PP_ALIGN.CENTER)
    _texto(s, x, 2.95, 2.95, 0.45, [(txt, 13, GRIS, False, False, 0)], align=PP_ALIGN.CENTER)
_rbox(s, 0.6, 3.75, 5.9, 2.35, fill=GRISCL, line=BORDE, line_w=1)
_box(s, 0.6, 3.75, 5.9, 0.5, fill=NARANJA)
_texto(s, 0.78, 3.8, 5.5, 0.42, [("EL PROBLEMA QUE RESUELVE", 14, BLANCO, True, False, 0)])
_texto(s, 0.8, 4.4, 5.5, 1.6, [
    ("Cada país regula la calidad del gas con SU PROPIA normativa, dispersa y en varios idiomas.",
     13, TINTA, False, True, 0),
    ("Las unidades y las condiciones NO coinciden: unos miden a 0 °C, otros a 15 o a 25 °C.",
     13, TINTA, False, True, 0),
    ("Compararlas a mano es lento y propenso a error.", 13, TINTA, False, True, 0),
], space=7)
_rbox(s, 6.85, 3.75, 5.9, 2.35, fill=GRISCL, line=BORDE, line_w=1)
_box(s, 6.85, 3.75, 5.9, 0.5, fill=VERDE)
_texto(s, 7.03, 3.8, 5.5, 0.42, [("CÓMO LO GARANTIZAMOS", 14, BLANCO, True, False, 0)])
_texto(s, 7.05, 4.4, 5.5, 1.6, [
    ("Las cifras NO nacen en el modelo: salen de una base verificada, con su cita.",
     13, AZUL, True, True, 0),
    ("Antes de comparar, se normalizan unidades y condiciones (ISO 13443).",
     13, TINTA, False, True, 0),
    ("Lo que la norma no fija, se declara: no se rellena con estimaciones.",
     13, TINTA, False, True, 0),
], space=7)
_box(s, 0.6, 6.25, 12.15, 0.65, fill=AZULCL)
_box(s, 0.6, 6.25, 0.09, 0.65, fill=AZUL)
_texto(s, 0.9, 6.27, 11.75, 0.6, [
    ("Todo lo anterior descansa en dos piezas: la ARQUITECTURA (que separa el dato del lenguaje) y la "
     "ONTOLOGÍA (donde vive el dato). Vamos a ellas.", 14, AZUL, True, False, 0),
], anchor=MSO_ANCHOR.MIDDLE)
_notes(s, "[~50 s] Recordad las cuatro cifras y saltad rápido al problema: la información regulatoria está "
          "dispersa y, sobre todo, NO ES DIRECTAMENTE COMPARABLE, porque cada país usa unidades y condiciones "
          f"distintas. Y la promesa: cero cifras inventadas. De los {TOT_CELDAS} registros trazables, "
          f"{TOT_V} están verificados verbatim, {TOT_VS} por fuente pública secundaria (cuando la norma "
          f"primaria es de pago) y {TOT_NV} declarados 'no verificable' porque la norma no los fija. "
          "Cerrad enlazando: 'esto se sostiene sobre dos piezas, la arquitectura y la ontología'.")

# ------------- 3. ARQUITECTURA — el esquema -------------
s = _slide(); _cabecera(s, "Arquitectura del sistema", nxt(), kicker="Pieza 1 · Arquitectura")
_texto(s, 0.6, 1.35, 12.2, 0.35,
       [("Cuatro componentes. La regla: toda cifra procede de la ontología; la IA solo redacta.",
         14, GRIS, False, False, 0)])
_rbox(s, 0.45, 3.05, 2.15, 1.85, fill=BLANCO, line=BORDE, line_w=1.25)
_box(s, 0.45, 3.05, 2.15, 0.42, fill=CYAN)
_texto(s, 0.45, 3.09, 2.15, 0.35, [("1 · INTERFAZ WEB", 11.5, BLANCO, True, False, 0)], align=PP_ALIGN.CENTER)
_texto(s, 0.6, 3.58, 1.9, 1.25, [
    ("El usuario pregunta", 13, AZUL, True, False, 0),
    ("index.html", 11, GRIS, False, False, 0),
    ("JavaScript vanilla", 11, GRIS, False, False, 0),
    ("marked · DOMPurify", 11, GRIS, False, False, 0),
], space=2)
_flecha(s, 2.63, 3.83, 0.32, 0.3, "der", CYAN)
_rbox(s, 2.98, 1.85, 5.2, 4.35, fill=GRISCL, line=AZUL, line_w=1.5)
_box(s, 2.98, 1.85, 5.2, 0.62, fill=AZUL)
_texto(s, 3.15, 1.9, 4.9, 0.52, [
    ("2 · SERVIDOR DE APLICACIÓN", 13, BLANCO, True, False, 0),
    ("Python · FastAPI · uvicorn · pydantic", 10.5, CYAN, False, False, 0),
], space=0)
_rbox(s, 3.2, 2.68, 4.75, 0.68, fill=AZUL2, line=None)
_texto(s, 3.2, 2.74, 4.75, 0.56, [
    ("ROUTER DETERMINISTA", 13, BLANCO, True, False, 0),
    ("¿la resuelve el código o hace falta la IA?", 10.5, RGBColor(0xBF, 0xDB, 0xEA), False, False, 0),
], align=PP_ALIGN.CENTER, space=0)
_flecha(s, 4.15, 3.4, 0.26, 0.24, "aba", VERDE)
_flecha(s, 6.55, 3.4, 0.26, 0.24, "aba", NARANJA)
_rbox(s, 3.2, 3.72, 4.75, 1.05, fill=VERDECL, line=VERDE, line_w=1.25)
_box(s, 3.2, 3.72, 0.09, 1.05, fill=VERDE)
_texto(s, 3.42, 3.79, 4.45, 0.92, [
    ("RUTA A · Consulta cuantitativa   (la mayoría)", 12.5, VERDEOS, True, False, 0),
    ("La resuelve el CÓDIGO. Sin IA.", 11.5, TINTA, False, False, 0),
    ("fuente_oficial · conversor_unidades · condiciones_referencia", 10, GRIS, False, False, 0),
], space=1)
_rbox(s, 3.2, 5.0, 4.75, 1.05, fill=NARANJACL, line=NARANJA, line_w=1.25)
_box(s, 3.2, 5.0, 0.09, 1.05, fill=NARANJA)
_texto(s, 3.42, 5.07, 4.45, 0.92, [
    ("RUTA B · Texto abierto", 12.5, NARANJAOS, True, False, 0),
    ("Va al modelo, que PIDE las cifras.", 11.5, TINTA, False, False, 0),
    ("llm_interface · function calling", 10, GRIS, False, False, 0),
], space=1)
_flecha(s, 8.2, 4.1, 0.34, 0.3, "der", VERDE)
_flecha(s, 8.2, 5.38, 0.34, 0.3, "der", NARANJA)
_rbox(s, 8.6, 3.35, 4.25, 1.5, fill=AZUL, line=None)
_texto(s, 8.8, 3.45, 3.9, 1.3, [
    ("3 · BASE DE CONOCIMIENTO", 11.5, CYAN, True, False, 0),
    ("Ontología YAML", 17, BLANCO, True, False, 0),
    (f"Los {TOT_CELDAS} registros, cada uno con su valor, unidad, condiciones y su cita oficial.",
     11, RGBColor(0xC9, 0xD6, 0xE0), False, False, 0),
], space=2)
_flecha(s, 9.05, 4.95, 0.3, 0.42, "arr", NARANJA)
_texto(s, 9.45, 4.98, 3.4, 0.4,
       [("function calling: el modelo PIDE la cifra, no la inventa", 10.5, NARANJAOS, True, False, 0)],
       anchor=MSO_ANCHOR.MIDDLE)
_rbox(s, 8.6, 5.45, 4.25, 1.0, fill=BLANCO, line=NARANJA, line_w=1.5, dash=True)
_texto(s, 8.8, 5.53, 3.9, 0.85, [
    ("4 · SERVICIO DE IA  (externo · opcional)", 11.5, NARANJAOS, True, False, 0),
    ("OpenAI GPT-4o-mini · temperatura 0 — solo REDACTA, nunca genera cifras.", 11, TINTA, False, False, 0),
], space=1)
_box(s, 0.45, 6.42, 12.4, 0.5, fill=AZULCL)
_box(s, 0.45, 6.42, 0.09, 0.5, fill=AZUL2)
_texto(s, 0.7, 6.46, 12.0, 0.42, [
    ("Si la IA no está disponible, el router conmuta a determinista: el chat nunca falla, y las cifras son las "
     "mismas.", 12, AZUL, False, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
_notes(s, "[~70 s] Recorred el esquema con el dedo. Uno: el usuario pregunta desde la interfaz. Dos: la "
          "pregunta entra en NUESTRO servidor y lo primero que encuentra es el ROUTER DETERMINISTA, que decide "
          "quién la resuelve. Si es cuantitativa —RUTA A, la mayoría— la resuelve el código leyendo la "
          "ontología: sin IA, sin posibilidad de inventar. Si es de texto abierto —RUTA B— va al modelo. Y "
          "aquí está lo importante: fijaos en la flecha que SUBE desde la IA hasta la ontología. Incluso en la "
          "ruta B el modelo no produce el número: lo PIDE, mediante function calling. Esa flecha es, "
          "técnicamente, la garantía anti-alucinación. Cerrad con la banda inferior: la IA es una comodidad, "
          "no una dependencia.")

# ------------- 4. ARQUITECTURA — las tres capas de datos -------------
s = _slide(); _cabecera(s, "Las tres capas de datos", nxt(), kicker="Pieza 1 · Arquitectura")
_texto(s, 0.6, 1.35, 12.2, 0.35,
       [("No hay una única base de datos: hay tres capas con funciones distintas. Y solo UNA guarda cifras.",
         14, GRIS, False, False, 0)])
capas = [
    ("CAPA 1", "Documentos oficiales", AZUL2,
     ["Los ~22 PDF de las normas: BOE, ERSE, GRTgaz, DVGW, Fluxys…",
      "Archivados en local para no depender de webs externas.",
      "Es la FUENTE ÚLTIMA DE VERDAD."], "PDF oficial"),
    ("CAPA 2", "La ontología", VERDE,
     [f"Los {TOT_CELDAS} registros extraídos de esos PDF, con su contexto y su cita.",
      "Un único fichero YAML, legible y versionado en git.",
      "De aquí salen TODAS las respuestas."], "PyYAML"),
    ("CAPA 3", "Índice documental (RAG)", NARANJA,
     ["pdfplumber trocea el TEXTO de los PDF en fragmentos con solape.",
      "Localiza el pasaje en las consultas de texto abierto.",
      "NO ALMACENA NINGUNA CIFRA."], "SQLite"),
]
x = 0.6
for etiqueta, tit, color, puntos, herr in capas:
    _rbox(s, x, 1.85, 3.93, 3.35, fill=GRISCL, line=BORDE, line_w=1)
    _box(s, x, 1.85, 3.93, 0.5, fill=color)
    _texto(s, x + 0.18, 1.9, 3.6, 0.42, [(etiqueta, 12, BLANCO, True, False, 0)])
    _texto(s, x + 0.18, 2.48, 3.6, 0.42, [(tit, 16, AZUL, True, False, 0)])
    _texto(s, x + 0.18, 2.98, 3.6, 1.7, [(p, 12, TINTA, False, True, 0) for p in puntos], space=7)
    _box(s, x + 0.18, 4.68, 1.6, 0.33, fill=BLANCO, line=color, line_w=1)
    _texto(s, x + 0.18, 4.7, 1.6, 0.29, [(herr, 11, color, True, False, 0)], align=PP_ALIGN.CENTER)
    x += 4.07
_box(s, 0.6, 5.35, 12.13, 0.62, fill=GRISCL)
_box(s, 0.6, 5.35, 0.09, 0.62, fill=CYAN)
_texto(s, 0.9, 5.38, 11.7, 0.56, [
    ("El stack completo", 11.5, CYAN, True, False, 0),
    ("Frontend: JavaScript vanilla · marked · DOMPurify    |    Backend: Python · FastAPI · uvicorn · pydantic"
     "    |    Datos: PyYAML · pdfplumber · SQLite    |    Informes: openpyxl · xhtml2pdf    |    "
     "IA (opcional): OpenAI GPT-4o-mini", 11, TINTA, False, False, 0),
], space=1)
_box(s, 0.6, 6.1, 12.13, 0.8, fill=AZULCL)
_box(s, 0.6, 6.1, 0.09, 0.8, fill=AZUL)
_texto(s, 0.9, 6.16, 11.7, 0.72, [
    ("Las CIFRAS residen solo en la capa 2. La capa 3 es un buscador de texto: no puede producir un número.",
     14, AZUL, True, False, 0),
    ("Todo el stack es software libre salvo la API de OpenAI — el único componente de pago, y además opcional.",
     12, GRIS, False, False, 0),
], space=2)
_notes(s, "[~50 s] Suele asumirse que hay una gran base de datos única. No la hay. Lo que hay que retener es "
          "el reparto: las cifras viven SOLO en la capa 2, la ontología. La capa 3, el índice documental, no "
          "guarda ningún número: solo localiza el pasaje pertinente cuando la consulta es de texto abierto. Y "
          "los PDF los guardamos en local para no depender de que una web externa siga viva. Si preguntan por "
          "pint, PyPDF2, pandas o cryptography: están en requirements pero HOY NO SE USAN, y lo sabemos.")

# ------------- 5. ONTOLOGÍA — la estructura real del fichero -------------
s = _slide(); _cabecera(s, "La ontología: estructura", nxt(), kicker="Pieza 2 · Ontología")
_texto(s, 0.6, 1.35, 12.2, 0.35,
       [("Un único fichero YAML. Esto es lo que hay dentro, exactamente.", 14, GRIS, False, False, 0)])
_rbox(s, 0.5, 1.8, 2.95, 3.85, fill=AZUL, line=None)
_texto(s, 0.7, 1.98, 2.6, 3.5, [
    ("FICHERO ÚNICO", 10.5, CYAN, True, False, 0),
    ("ontologia_enagas.yaml", 14, BLANCO, True, False, 0),
    (f"v{VER_ONT} · rev. {REV_ONT}", 10.5, RGBColor(0x9F, 0xB3, 0xC4), False, False, 0),
    ("", 5, BLANCO, False, False, 0),
    ("Legible por una persona.", 11.5, RGBColor(0xC9, 0xD6, 0xE0), False, False, 0),
    ("Versionado en git junto al código: cada cambio de una cifra queda registrado.",
     11.5, RGBColor(0xC9, 0xD6, 0xE0), False, False, 0),
    ("", 5, BLANCO, False, False, 0),
    ("Por eso YAML y no una base de datos: a esta escala es más AUDITABLE.",
     11.5, CYAN, False, False, 0),
], space=3)
_box(s, 3.78, 1.95, 0.045, 3.5, fill=BORDE)
ramas = [
    ("ontologia.fuentes_normativas", f"El catálogo de las {N_FUENTES} normas oficiales: organismo, publicación, URL y copia local del PDF.", AZUL2),
    ("ontologia.tipos_gas", f"Los {N_GASES} tipos de gas y a qué sección de parámetros apunta cada uno.", AZUL2),
    ("unidades", "El catálogo de unidades admitidas (% mol, mg/m3, kWh/m3, MJ/m3, °C…).", AZUL2),
    ("parametros", f"GAS NATURAL — {GN_P} parámetros × {GN_J} jurisdicciones = {GN_C} celdas.", VERDE),
    ("parametros_biometano", f"BIOMETANO — {BIO_P} parámetros · {BIO_C} celdas ({BIO_J} jurisdicciones: ES · PT · FR · UE).", VERDE),
    ("parametros_hidrogeno", f"HIDRÓGENO — {H2_P} parámetros · {H2_C} celdas (dominio de red + producto).", VERDE),
]
y = 1.85
for clave, desc, color in ramas:
    _box(s, 3.825, y + 0.24, 0.4, 0.045, fill=BORDE)
    _rbox(s, 4.3, y, 8.55, 0.53, fill=GRISCL, line=BORDE, line_w=1)
    _box(s, 4.3, y, 0.08, 0.53, fill=color)
    _texto(s, 4.48, y + 0.02, 3.15, 0.49, [(clave, 11.5, color, True, False, 0)],
           anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    _texto(s, 7.7, y + 0.02, 5.0, 0.49, [(desc, 11, TINTA, False, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.62
_box(s, 4.3, 5.65, 8.55, 0.55, fill=NARANJACL)
_box(s, 4.3, 5.65, 0.08, 0.55, fill=NARANJA)
_texto(s, 4.48, 5.67, 8.2, 0.51, [
    (f"{' · '.join(CLAVES_INERTES)}  —  documentación de diseño dentro del YAML: el código NO las lee.",
     11, NARANJAOS, True, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
_box(s, 0.5, 6.35, 12.35, 0.55, fill=AZULCL)
_box(s, 0.5, 6.35, 0.09, 0.55, fill=AZUL)
_texto(s, 0.8, 6.37, 11.9, 0.51, [
    ("Las tres secciones de parámetros comparten EXACTAMENTE el mismo esquema: por eso ampliar a biometano e "
     "hidrógeno no obligó a cambiar la arquitectura.", 13, AZUL, True, False, 0),
], anchor=MSO_ANCHOR.MIDDLE)
_notes(s, "[~70 s] Esta es la estructura REAL del fichero, no una versión idealizada. Arriba, lo que el código "
          f"lee de verdad: el catálogo de las {N_FUENTES} normas oficiales, el registro de los {N_GASES} tipos "
          "de gas, el catálogo de unidades, y las tres secciones de parámetros —una por gas—. Lo importante: "
          "esas tres secciones comparten EXACTAMENTE el mismo esquema, y por eso ampliar a biometano e "
          "hidrógeno no obligó a tocar la arquitectura. Y sed honestos con la caja naranja: hay claves "
          f"({', '.join(CLAVES_INERTES)}) que son documentación del diseño dentro del propio YAML y que el "
          "código NO lee. Decirlo vosotros primero vale más que si os lo encuentran. Por último, el porqué del "
          "YAML: a esta escala un fichero legible y versionado en git es más auditable que una base de datos.")

# ------------- 6. ONTOLOGÍA — anatomía de un valor y la garantía -------------
s = _slide(); _cabecera(s, "Anatomía de un valor, y la garantía", nxt(), kicker="Pieza 2 · Ontología")
_texto(s, 0.6, 1.35, 12.2, 0.35,
       [("Cada límite NO es solo un número: guarda todo su contexto normativo. Ejemplo real — el O2 de España:",
         14, GRIS, False, False, 0)])
_rbox(s, 0.55, 1.8, 6.35, 3.75, fill=RGBColor(0xF5, 0xF8, 0xFA), line=BORDE, line_w=1.25)
_box(s, 0.55, 1.8, 6.35, 0.4, fill=AZUL2)
_texto(s, 0.72, 1.83, 6.0, 0.34,
       [("parametros -> O2 -> limites -> ES", 10.5, BLANCO, False, False, 0)], font=MONO)
_texto(s, 0.78, 2.3, 6.0, 3.15, [
    ("ES:", 12, AZUL, True, False, 0),
    ("  fuente: ORDEN_TED_181_2025", 11.5, TINTA, False, False, 0),
    ('  articulo: "Tabla 3, apdo. 2.5.2.1 (pág. 27)"', 11.5, TINTA, False, False, 0),
    ("  tipo_limite: maximo", 11.5, TINTA, False, False, 0),
    ("  valor: 0.01", 11.5, TINTA, False, False, 0),
    ("  unidad: pct_mol", 11.5, TINTA, False, False, 0),
    ('  expresion_original: "O2: – / 0,01 % mol"', 11.5, AZUL, True, False, 0),
    ("  condiciones_referencia:", 11.5, TINTA, False, False, 0),
    ("    temperatura_volumen_C: 0", 11.5, TINTA, False, False, 0),
    ("    presion_bar: 1.01325", 11.5, TINTA, False, False, 0),
    ("  estado_verificacion: VERIFICADO", 11.5, VERDEOS, True, False, 0),
], space=3, font=MONO)
_texto(s, 0.78, 5.05, 6.0, 0.45, [
    ("La cita (fuente + artículo) y el TEXTO LITERAL permiten recomprobar la celda sin abrir el PDF. "
     "Y condiciones_referencia es lo que hace posible normalizar con ISO 13443 antes de comparar.",
     10.5, GRIS, False, False, 0)])
# Estados de verificación (los tres reales)
estados = [
    (TOT_V, "VERIFICADO", "Contrastado VERBATIM contra el boletín oficial.", VERDECL, VERDE, VERDEOS),
    (TOT_VS, "VERIFICADO_SECUNDARIO", "La norma primaria es de pago: valor tomado de fuente pública secundaria, citada.",
     AZULCL, CYAN, AZUL),
    (TOT_NV, "NO VERIFICABLE", "La norma NO FIJA ese parámetro. No se inventa: se declara el hueco.",
     NARANJACL, NARANJA, NARANJAOS),
]
y = 1.8
for n, tit, desc, bg, barra, txtcol in estados:
    _rbox(s, 7.15, y, 5.67, 1.15, fill=bg, line=barra, line_w=1.25)
    _box(s, 7.15, y, 0.09, 1.15, fill=barra)
    _texto(s, 7.42, y + 0.06, 1.35, 0.55, [(str(n), 26, txtcol, True, False, 0)])
    _texto(s, 8.75, y + 0.05, 3.9, 1.05, [
        (tit, 12, txtcol, True, False, 0),
        (desc, 10.5, TINTA, False, False, 0),
    ], space=1)
    y += 1.3
_box(s, 7.15, 5.7, 5.67, 0.5, fill=GRISCL)
_texto(s, 7.42, 5.72, 5.2, 0.46,
       [(f"{TOT_CELDAS} celdas en total  ·  gas natural: {N_VERIF} verificadas / {N_NOVER} no verificables",
         11, AZUL, True, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
_box(s, 0.55, 6.35, 12.27, 0.55, fill=AZULCL)
_box(s, 0.55, 6.35, 0.09, 0.55, fill=AZUL)
_texto(s, 0.85, 6.37, 11.85, 0.51, [
    ("DINAMARCA: los límites de O2/CO2 de su norma son del BIOGÁS DE DISTRIBUCIÓN, no del gas natural de "
     "transporte. Por eso se marcaron NO VERIFICABLE en vez de rellenar el hueco.", 12.5, AZUL, True, False, 0),
], anchor=MSO_ANCHOR.MIDDLE)
_notes(s, "[~60 s] Aquí está el valor real del diseño: de cada dato no guardamos solo el número, sino todo su "
          "contexto. Señalad dos campos. Uno, 'expresion_original': el texto literal de la norma, tal cual está "
          "redactado —eso convierte la confianza en VERIFICABILIDAD: no hay que creerse el dato, se comprueba—. "
          "Dos, 'condiciones_referencia': es lo que permite normalizar con la ISO 13443 antes de comparar, "
          "porque comparar en bruto valores medidos a 0, 15 o 25 grados sería metodológicamente incorrecto. A "
          "la derecha, los tres estados. Y el ejemplo de Dinamarca es el que mejor lo ilustra: era fácil "
          "rellenar el hueco con un número que estaba en la misma norma, pero pertenecía a otro contexto "
          "—biogás de distribución, no transporte—. Preferimos el hueco honesto: informa mejor que un número "
          "falso.")

# ------------- 7. Cierre — las garantías -------------
s = _slide()
_box(s, 0, 0, 13.333, 7.5, fill=AZUL)
_box(s, 0, 1.45, 13.333, 0.05, fill=CYAN)
_texto(s, 0.9, 0.5, 11.5, 0.8, [("Las garantías del sistema", 32, BLANCO, True, False, 0)])
gar = [
    ("Cero cifras inventadas", "Los valores proceden de código y datos verificados, nunca del modelo."),
    ("Trazabilidad completa", "Cada valor cita su norma, su artículo, su página y su enlace."),
    ("Transparencia", "Lo que la norma no fija se declara explícitamente; no se completa."),
    ("Reproducibilidad", "Con el mismo código y los mismos datos, el resultado es idéntico en cualquier entorno."),
    ("Auditabilidad", "La base de conocimiento es consultable: se puede verificar el origen de cada valor."),
]
y = 2.0
for t, d in gar:
    _box(s, 0.9, y, 0.07, 0.78, fill=VERDE)
    _texto(s, 1.25, y - 0.03, 11.3, 0.42, [(t, 18, CYAN, True, False, 0)])
    _texto(s, 1.25, y + 0.33, 11.3, 0.42, [(d, 14, RGBColor(0xC9, 0xD6, 0xE0), False, False, 0)])
    y += 0.92
_box(s, 0.9, 6.6, 11.6, 0.06, fill=VERDE)
_texto(s, 0.9, 6.75, 11.6, 0.5,
       [("La IA redacta el texto. Los números salen siempre de la normativa oficial.",
         18, BLANCO, True, False, 0)])
nxt()
_notes(s, "[~30 s] Cierre. Las cinco garantías salen del rigor del proceso, no de la suerte. Y la frase final "
          "es la que queremos que se lleven: la IA redacta el texto; los números salen siempre de la normativa "
          "oficial. Gracias, y quedamos a vuestra disposición para las preguntas.")

prs.save(PPTX)
print("PPTX generado:", os.path.relpath(PPTX, RAIZ),
      f"({os.path.getsize(PPTX)//1024} KB, {len(prs.slides._sldIdLst)} diapositivas)")
