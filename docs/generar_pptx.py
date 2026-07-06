# -*- coding: utf-8 -*-
"""Genera el PowerPoint de presentación del Comparador (estilo Enagás).

    python docs/generar_pptx.py

Diapositivas limpias (bullets, tablas, un diagrama de arquitectura) y NOTAS DEL PONENTE
con el guion detallado en cada slide, para que la explicación esté completa.
Salida: docs/Presentacion_Comparador_Gas.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

AQUI = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(AQUI, "Presentacion_Comparador_Gas.pptx")

# --- Paleta Enagás ---
AZUL   = RGBColor(0x01, 0x3A, 0x57)   # azul corporativo oscuro
CYAN   = RGBColor(0x00, 0x99, 0xD6)   # cian de acento
VERDE  = RGBColor(0x6C, 0xB3, 0x3E)   # verde
GRIS   = RGBColor(0x4A, 0x5B, 0x68)   # texto secundario
TINTA  = RGBColor(0x1B, 0x2A, 0x38)   # texto principal
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRISCL = RGBColor(0xEE, 0xF2, 0xF6)   # fondo suave
FUENTE = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def _slide():
    return prs.slides.add_slide(BLANK)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _box(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """runs: lista de párrafos; cada uno es (texto, size, color, bold, bullet, level)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, r in enumerate(runs):
        txt, size, color, bold, bullet, level = (r + (None,)*6)[:6]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.level = level or 0
        run = p.add_run(); run.text = ("• " + txt) if bullet else txt
        run.font.name = FUENTE; run.font.size = Pt(size)
        run.font.color.rgb = color; run.font.bold = bool(bold)
    return tb


def _cabecera(slide, titulo, n):
    # barra superior
    _box(slide, 0, 0, W, Inches(1.05), fill=AZUL)
    _box(slide, 0, Inches(1.05), W, Pt(4), fill=CYAN)
    _text(slide, Inches(0.55), Inches(0.16), Inches(11.5), Inches(0.8),
          [(titulo, 26, BLANCO, True, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
    # pie
    _text(slide, Inches(0.55), Inches(7.05), Inches(9), Inches(0.35),
          [("Comparador Regulatorio de Calidad de Gas Natural — Enagás", 10, GRIS, False, False, 0)])
    _text(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
          [(str(n), 10, GRIS, False, False, 0)], align=PP_ALIGN.RIGHT)


def slide_contenido(titulo, n, intro=None, bullets=None, notas="", bcolor=TINTA, bsize=17):
    s = _slide(); _cabecera(s, titulo, n)
    top = 1.45
    if intro:
        _text(s, Inches(0.6), Inches(top), Inches(12.1), Inches(0.9),
              [(intro, 15, GRIS, False, False, 0)])
        top += 0.85
    if bullets:
        runs = []
        for b in bullets:
            if isinstance(b, tuple):  # (texto, sub?)  sub=True -> nivel 2
                runs.append((b[0], bsize-2 if b[1] else bsize, GRIS if b[1] else bcolor, False, True, 1 if b[1] else 0))
            else:
                runs.append((b, bsize, bcolor, False, True, 0))
        _text(s, Inches(0.7), Inches(top), Inches(11.9), Inches(5.2), runs, space=10)
    if notas:
        _notes(s, notas)
    return s


def slide_tabla(titulo, n, headers, filas, notas="", intro=None, col_ratios=None):
    s = _slide(); _cabecera(s, titulo, n)
    top = 1.5
    if intro:
        _text(s, Inches(0.6), Inches(top), Inches(12.1), Inches(0.8),
              [(intro, 15, GRIS, False, False, 0)])
        top += 0.75
    nrows, ncols = len(filas) + 1, len(headers)
    tw = Inches(12.1)
    gr = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(top), tw, Inches(0.5)).table
    gr.first_row = False; gr.horz_banding = False
    if col_ratios:
        tot = sum(col_ratios)
        for j, cr in enumerate(col_ratios):
            gr.columns[j].width = int(tw * cr / tot)
    # cabecera
    for j, htxt in enumerate(headers):
        c = gr.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = AZUL
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = htxt
        r.font.name = FUENTE; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = BLANCO
    # filas
    for i, fila in enumerate(filas, start=1):
        for j, val in enumerate(fila):
            c = gr.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = GRISCL if i % 2 else BLANCO
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = Pt(3); c.margin_bottom = Pt(3)
            p = c.text_frame.paragraphs[0]; p.word_wrap = True
            for k, seg in enumerate((val if isinstance(val, list) else [(val, False)])):
                txt, bold = seg
                run = p.add_run(); run.text = txt
                run.font.name = FUENTE; run.font.size = Pt(13)
                run.font.bold = bool(bold); run.font.color.rgb = TINTA
    if notas:
        _notes(s, notas)
    return s


# =========================== CONTENIDO ===========================
N = 0
def nxt():
    global N; N += 1; return N


# --- 1. Portada ---
s = _slide()
_box(s, 0, 0, W, H, fill=AZUL)
_box(s, 0, Inches(4.55), W, Pt(5), fill=CYAN)
_box(s, Inches(0.9), Inches(2.0), Inches(0.16), Inches(2.2), fill=VERDE)
_text(s, Inches(1.25), Inches(2.0), Inches(11), Inches(2.3), [
    ("Comparador Regulatorio de", 40, BLANCO, True, False, 0),
    ("Calidad de Gas Natural", 40, BLANCO, True, False, 0),
], space=2)
_text(s, Inches(1.28), Inches(4.75), Inches(11), Inches(1.2), [
    ("Arquitectura y funcionamiento del sistema", 20, CYAN, False, False, 0),
    ("Comparativa regulatoria entre 21 jurisdicciones y 10 parámetros de calidad", 15, RGBColor(0xC9,0xD6,0xE0), False, False, 0),
], space=6)
_text(s, Inches(1.28), Inches(6.5), Inches(11), Inches(0.5),
      [("Documento de presentación", 13, RGBColor(0x9F,0xB3,0xC4), False, False, 0)])
_notes(s, "Presentación del sistema ante Enagás. Objetivo: explicar de forma clara qué hace la aplicación, "
          "cómo está construida y por qué es fiable. El hilo de la exposición es por componentes y capas.")

# --- 2. Índice ---
slide_contenido("Índice", nxt(),
    intro="Recorrido de la presentación:",
    bullets=[
        "Contexto y objetivo",
        "Arquitectura general del sistema",
        "Las tres capas de datos y la base de conocimiento",
        "El servidor de aplicación (FastAPI) frente a la API de OpenAI",
        "El motor determinista y la normalización (ISO 13443)",
        "La capa de inteligencia artificial y la recuperación documental (RAG)",
        "Metodología de construcción, verificación y garantías",
    ],
    notas="Presentar la agenda. Insistir en que el hilo es por componentes: primero la visión general, "
          "luego cada pieza, y por último la metodología y las garantías.")

# --- 3. Contexto y objetivo ---
slide_contenido("Contexto y objetivo", nxt(),
    intro="El problema y la propuesta de valor.",
    bullets=[
        "Cada país regula la calidad admisible del gas natural con su propia normativa (poder calorífico, Wobbe, azufre, CO₂, puntos de rocío…).",
        "Están dispersas en boletines oficiales distintos, en varios idiomas y con unidades y condiciones de referencia diferentes.",
        "Comparar dos marcos a mano es laborioso y propenso a error.",
        "Solución: un asistente que compara esa calidad entre 21 jurisdicciones y 10 parámetros, en lenguaje natural.",
        "Principio de diseño: AUSENCIA DE CIFRAS NO VERIFICADAS. Ningún valor se estima; todos proceden de normativa oficial, con su cita.",
    ],
    notas="Cada país regula la calidad admisible del gas mediante su propia normativa. Esa información está "
          "dispersa, en varios idiomas y con unidades distintas; compararla a mano es costoso y arriesgado. "
          "Hemos desarrollado un sistema que la compara entre 21 jurisdicciones y 10 parámetros. Su principio "
          "de diseño es que no genera ninguna cifra por estimación: todas proceden de una base contrastada "
          "frente a la normativa oficial y se presentan con su cita.")

# --- 4. Cifras clave ---
s = _slide(); _cabecera(s, "El sistema en cifras", nxt())
tarjetas = [("21", "jurisdicciones"), ("10", "parámetros de calidad"),
            ("210", "valores verificables"), ("0", "cifras inventadas")]
x = 0.6
for num, txt in tarjetas:
    _box(s, Inches(x), Inches(2.3), Inches(2.95), Inches(2.6), fill=GRISCL)
    _box(s, Inches(x), Inches(2.3), Inches(2.95), Pt(6), fill=CYAN)
    _text(s, Inches(x), Inches(2.75), Inches(2.95), Inches(1.3),
          [(num, 54, AZUL, True, False, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(x), Inches(4.0), Inches(2.95), Inches(0.7),
          [(txt, 15, GRIS, False, False, 0)], align=PP_ALIGN.CENTER)
    x += 3.08
_text(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.8),
      [("Cobertura: 175 valores VERIFICADOS y 35 NO VERIFICABLES (la norma no los fija) — nunca inventados.",
        15, TINTA, True, False, 0)], align=PP_ALIGN.CENTER)
_notes(s, "Cuatro cifras que resumen el alcance: 21 jurisdicciones, 10 parámetros, 210 valores y cero cifras "
          "inventadas. De los 210, 175 están verificados verbatim contra su boletín oficial y 35 se marcan como "
          "'no verificable' porque la norma de ese país no los fija.")

# --- 5. Arquitectura general (diagrama) ---
s = _slide(); _cabecera(s, "Arquitectura general", nxt())
_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.6),
      [("Cuatro componentes con responsabilidades bien delimitadas.", 15, GRIS, False, False, 0)])
def _comp(l, t, w, h, tit, sub, fill, tcol=BLANCO):
    b = _box(s, Inches(l), Inches(t), Inches(w), Inches(h), fill=fill)
    _text(s, Inches(l), Inches(t+0.12), Inches(w), Inches(h-0.2),
          [(tit, 16, tcol, True, False, 0), (sub, 11.5, tcol, False, False, 0)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=3)
    return b
def _flecha(l, t, w, h):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = CYAN; a.line.fill.background(); a.shadow.inherit = False
_comp(4.9, 2.1, 3.5, 0.85, "INTERFAZ WEB", "El usuario formula su consulta", CYAN)
_flecha(6.45, 3.02, 0.4, 0.5)
_comp(4.4, 3.6, 4.5, 0.95, "SERVIDOR DE APLICACIÓN", "FastAPI · recibe, aplica la lógica y orquesta", AZUL)
_flecha(3.7, 4.65, 0.35, 0.5); _flecha(9.25, 4.65, 0.35, 0.5)
_comp(1.35, 5.25, 4.3, 1.15, "BASE DE CONOCIMIENTO", "Ontología: 210 cifras verificadas + sus fuentes", VERDE)
_comp(7.65, 5.25, 4.3, 1.15, "SERVICIO DE IA (externo)", "OpenAI: solo texto abierto · no genera cifras", GRIS)
_notes(s, "El sistema se estructura en cuatro componentes. La interfaz web presenta la información al usuario. "
          "El servidor de aplicación, desarrollado con FastAPI, concentra la lógica y decide cómo resolver cada "
          "consulta. La base de conocimiento (la ontología) almacena los datos verificados. Y el servicio de IA "
          "es un proveedor externo que se emplea de forma acotada, solo para redactar respuestas de texto, sin "
          "capacidad de generar cifras.")

# --- 6. Las tres capas de datos ---
slide_tabla("Las tres capas de datos", nxt(),
    ["Capa", "Contenido", "Función"],
    [
        [[("1. Documentos oficiales", True)], "PDF de las normas (BOE, ERSE, DVGW, Fluxys, National Grid…)", "Fuente primaria y última de verdad"],
        [[("2. Ontología", True)], "Fichero estructurado con las 210 cifras extraídas de esos PDF", "Repositorio del que salen las respuestas"],
        [[("3. Índice documental (RAG)", True)], "Índice del texto de los PDF, segmentado por página", "Buscador interno para consultas abiertas"],
    ],
    intro="No hay una única base de datos gigante: la información se organiza en tres capas.",
    col_ratios=[2.2, 4.2, 3.2],
    notas="Es habitual asumir que hay una gran base de datos única; no es el caso. Las cifras residen en la "
          "capa 2, la ontología, y cada una referencia su documento oficial de la capa 1. La capa 3 no almacena "
          "ninguna cifra: es un índice de búsqueda sobre el texto de los documentos, para las consultas abiertas.")

# --- 7. La base de conocimiento (ontología) ---
slide_contenido("La base de conocimiento (ontología)", nxt(),
    intro="El elemento central. De cada valor se registra su valor y todo su contexto normativo:",
    bullets=[
        "El valor numérico (o su rango) y la unidad.",
        "Las condiciones de referencia (a qué temperatura se mide).",
        "El texto literal de la norma, tal como está redactado.",
        "La cita completa: norma, artículo, página y enlace.",
        "Una nota aclaratoria con los matices.",
        "El estado de verificación (garantía frente a la invención de datos).",
    ],
    notas="La ontología es el elemento central. Para cada jurisdicción y parámetro se registra no solo el valor, "
          "sino todo su contexto: unidad, condiciones de referencia, el texto literal de la norma, la cita completa "
          "(norma, artículo, página y enlace), una nota aclaratoria y el estado de verificación. "
          "Nota técnica: se usa un fichero estructurado (YAML) en lugar de una base de datos relacional porque, "
          "a esta escala, es más auditable y trazable, y se versiona junto con el código.")

# --- 8. Estados de verificación ---
s = _slide(); _cabecera(s, "Estados de verificación", nxt())
_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.6),
      [("La garantía frente a la invención de datos: cada cifra está en uno de dos estados.", 15, GRIS, False, False, 0)])
_box(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(2.7), fill=GRISCL); _box(s, Inches(0.7), Inches(2.2), Inches(5.8), Pt(6), fill=VERDE)
_text(s, Inches(1.0), Inches(2.45), Inches(5.2), Inches(2.4), [
    ("VERIFICADO", 22, AZUL, True, False, 0),
    ("175 valores", 15, VERDE, True, False, 0),
    ("Cifra contrastada literalmente contra su boletín oficial.", 14, TINTA, False, False, 0),
], space=8)
_box(s, Inches(6.85), Inches(2.2), Inches(5.8), Inches(2.7), fill=GRISCL); _box(s, Inches(6.85), Inches(2.2), Inches(5.8), Pt(6), fill=CYAN)
_text(s, Inches(7.15), Inches(2.45), Inches(5.2), Inches(2.4), [
    ("NO VERIFICABLE", 22, AZUL, True, False, 0),
    ("35 valores", 15, CYAN, True, False, 0),
    ("La norma de esa jurisdicción no fija ese parámetro. No se estima: se marca y se explica.", 14, TINTA, False, False, 0),
], space=8)
_text(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.0),
      [("No existe un estado intermedio: un valor consta en la norma, o se declara explícitamente que la norma no lo establece.",
        15, TINTA, True, False, 0)], align=PP_ALIGN.CENTER)
_notes(s, "El estado de verificación es la garantía anti-invención. Verificado: 175 valores contrastados verbatim. "
          "No verificable: 35 valores que la norma del país no fija; no se completan con una estimación, se marcan "
          "y se explica el motivo. No hay un tercer estado intermedio. Ejemplo: en Dinamarca, los límites de O₂ y "
          "CO₂ de la norma son de biogás, no de gas natural, así que para gas natural se dejaron como no verificable.")

# --- 9. FastAPI vs API de OpenAI ---
slide_tabla("FastAPI y la API de OpenAI: son cosas distintas", nxt(),
    ["", "FastAPI", "API de OpenAI"],
    [
        [[("Naturaleza", True)], "Framework para construir NUESTRO servidor", "Servicio externo que consumimos"],
        [[("Titularidad", True)], "Propia (es nuestro backend)", "De OpenAI (somos cliente)"],
        [[("Coste", True)], "Sin coste (código abierto)", "De pago, por uso"],
        [[("Papel", True)], "Núcleo de la aplicación", "Proveedor auxiliar, invocado de forma controlada"],
    ],
    intro="Ambos incluyen el término «API», pero son componentes de niveles distintos.",
    col_ratios=[1.5, 4.0, 4.0],
    notas="Conviene distinguir dos componentes que a veces se confunden por su nombre. FastAPI es el framework "
          "con el que hemos desarrollado nuestro servidor; es infraestructura propia y gratuita. La API de OpenAI "
          "es un servicio de terceros que consumimos puntualmente para tareas de lenguaje. El servidor es "
          "imprescindible: atiende la web, accede a los datos, ejecuta los cálculos exactos, decide cuándo usar la "
          "IA (en la mayoría de casos no la usa) y custodia las credenciales, que nunca se exponen en el navegador.")

# --- 10. El motor determinista ---
slide_contenido("El motor determinista: enrutado de consultas", nxt(),
    intro="Toda consulta pasa primero por un enrutado que decide cómo resolverla. «Determinista» = misma consulta, misma respuesta, por código, sin IA.",
    bullets=[
        "Consultas CUANTITATIVAS (un límite, un cumplimiento, una comparación, una conversión): las resuelve el código leyendo la ontología. Sin IA, sin posibilidad de generar un valor incorrecto.",
        "Consultas de TEXTO ABIERTO («¿en qué consiste el índice de Wobbe?»): se derivan al servicio de IA.",
        ("En la práctica, la mayoría de consultas se resuelven sin recurrir a la IA.", True),
    ],
    notas="Toda consulta pasa primero por el motor determinista. Determinista significa que, ante la misma "
          "consulta, produce siempre la misma respuesta, calculada por código, sin aleatoriedad ni IA. Distingue "
          "dos tipos: las cuantitativas (límites, cumplimiento, comparaciones, conversiones), que resuelve el "
          "código leyendo la ontología, sin IA; y las de texto abierto, que deriva al servicio de IA. La mayoría "
          "se resuelven sin IA.")

# --- 11. Normalización ISO 13443 ---
slide_contenido("Normalización de condiciones (ISO 13443)", nxt(),
    intro="Comparar de forma rigurosa: cada país usa unidades y condiciones de referencia distintas.",
    bullets=[
        "Unos expresan en kWh/m³, otros en MJ/m³; unos miden a 0 °C, otros a 15 o 25 °C.",
        "Comparar los valores en bruto sería metodológicamente incorrecto (como comparar millas con kilómetros).",
        "Todos los valores se llevan a la base de referencia española con los factores de la norma ISO 13443.",
        "Esos factores no se estiman: se toman literalmente de la norma y están implementados y verificados.",
        ("Los valores derivados se muestran con 2 decimales (sin falsa precisión); los originales, con su precisión de origen.", True),
    ],
    notas="Un aspecto clave para la credibilidad es la comparabilidad. Cada país expresa sus límites en unidades "
          "y condiciones distintas; compararlos en bruto sería incorrecto. Para una comparación rigurosa, todos "
          "los valores se llevan a la base española aplicando los factores de la norma internacional ISO 13443, "
          "que se toman literalmente y están verificados. Así, al comparar dos jurisdicciones, ambas cifras están "
          "en la misma base.")

# --- 12. La IA y sus salvaguardas ---
slide_contenido("La inteligencia artificial y sus salvaguardas", nxt(),
    intro="Cuando una consulta se deriva al servicio de IA, este opera bajo restricciones estrictas:",
    bullets=[
        "Tiene PROHIBIDO generar cifras: si necesita un dato, lo solicita a las herramientas internas, que lo obtienen de la ontología.",
        "Debe CITAR los documentos oficiales.",
        "Su ámbito se limita a la calidad del gas natural.",
        "Si el servicio de IA no está disponible, el sistema conmuta automáticamente al modo determinista: nunca se interrumpe.",
        ("Modelo: OpenAI GPT-4o-mini, con temperatura 0 (máxima previsibilidad).", True),
    ],
    notas="Cuando una consulta llega al servicio de IA, este opera atado en corto: tiene prohibido generar cifras "
          "(si necesita un número, lo pide a las herramientas internas, que lo sacan de la ontología), debe citar "
          "los documentos y se limita a la calidad del gas. Además, si la IA no está disponible, el sistema "
          "conmuta automáticamente al modo determinista, de modo que el servicio nunca se interrumpe.")

# --- 13. RAG ---
slide_contenido("Recuperación documental (RAG)", nxt(),
    intro="Para las consultas de texto abierto, la respuesta se fundamenta en los documentos oficiales, no en el conocimiento general del modelo.",
    bullets=[
        "Fase 1 — INDEXACIÓN: el sistema procesa los PDF, extrae su texto, lo segmenta por página y lo almacena en un índice de búsqueda.",
        "Fase 2 — RECUPERACIÓN: cuando la IA necesita fundamentar una respuesta, localiza los fragmentos pertinentes (con documento y página) y redacta citándolos.",
        ("La búsqueda es léxica (por términos), no semántica: más simple, suficiente y plenamente reproducible.", True),
        ("La indexación es incremental: solo se reprocesa lo nuevo o modificado → arranque casi inmediato.", True),
    ],
    notas="El sistema incorpora recuperación documental (RAG). Su fin es que, en consultas de texto abierto, la "
          "respuesta se apoye en los documentos oficiales. Funciona en dos fases: indexación (procesar los PDF y "
          "almacenar su texto segmentado por página) y recuperación (localizar los fragmentos pertinentes y "
          "redactar citándolos). La búsqueda es léxica, no semántica: es reproducible, no depende de una caja "
          "negra externa.")

# --- 14. Metodología ---
slide_contenido("Metodología de construcción y verificación", nxt(),
    intro="La fiabilidad es consecuencia del rigor del proceso de carga. Para cada jurisdicción:",
    bullets=[
        "1. Se identificó la normativa oficial vigente.",
        "2. Se obtuvo el documento y se archivó localmente.",
        "3. Se transcribió cada cifra literalmente, sin interpretación.",
        "4. Se verificó individualmente contra el documento.",
        "5. Lo que la norma no establece, no se completa: se marca como «no verificable», con su justificación.",
        "6. Se incorporó la normalización (ISO 13443) para comparaciones homogéneas.",
        ("Controles automáticos: comprueban que los 210 valores se resuelven, que los enlaces funcionan y que no hay incoherencias.", True),
    ],
    notas="La fiabilidad es consecuencia del rigor del proceso de carga: identificar la norma vigente, obtener el "
          "documento, transcribir cada cifra literalmente, verificarla una a una, marcar como no verificable lo "
          "que la norma no fija, y normalizar con la ISO 13443. Además hay controles de calidad automatizados.",
    bsize=16)

# --- 15. Garantías ---
slide_contenido("Garantías del sistema", nxt(),
    bullets=[
        "Ausencia de cifras inventadas, por diseño: los valores proceden de código y datos verificados, no del modelo.",
        "Trazabilidad completa: cada valor cita norma, artículo, página y enlace.",
        "Transparencia: lo que la norma no fija se declara explícitamente; no se completa.",
        "Reproducibilidad: con el mismo código y datos, el resultado es idéntico en cualquier entorno.",
        "Auditabilidad: la base de conocimiento es consultable y permite verificar el origen de cada valor.",
    ],
    notas="Cinco garantías que se llevan: cero cifras inventadas por diseño, trazabilidad completa (cada valor "
          "con su cita), transparencia (lo que no está en la norma se dice), reproducibilidad (mismo resultado en "
          "cualquier entorno) y auditabilidad (se puede verificar el origen de cada valor).",
    bcolor=TINTA, bsize=18)

# --- 16. Historial persistente ---
slide_contenido("El chat: historial persistente", nxt(),
    intro="La conversación no se pierde.",
    bullets=[
        "Cada consulta y su respuesta se guardan en el navegador y se restauran al recargar la página o tras reiniciar el servidor.",
        "El usuario no pierde su historial de consultas.",
        "La opción «Nueva consulta» inicia una sesión limpia.",
        "El asistente conserva el contexto de la conversación para las preguntas de seguimiento.",
    ],
    notas="La conversación es persistente: cada pregunta y respuesta se guarda en el navegador y se restaura al "
          "recargar o tras reiniciar el servidor, de modo que el usuario no pierde su historial. «Nueva consulta» "
          "empieza una sesión limpia, y el asistente conserva el contexto para las preguntas de seguimiento.")

# --- 17. Cierre ---
s = _slide()
_box(s, 0, 0, W, H, fill=AZUL)
_box(s, Inches(0.9), Inches(2.4), Inches(0.16), Inches(2.6), fill=VERDE)
_text(s, Inches(1.3), Inches(2.5), Inches(11), Inches(2.6), [
    ("En síntesis", 22, CYAN, True, False, 0),
    ("Un asistente que compara la calidad regulatoria del gas natural", 26, BLANCO, True, False, 0),
    ("entre 21 jurisdicciones, con trazabilidad total y sin cifras inventadas.", 26, BLANCO, True, False, 0),
], space=8)
_text(s, Inches(1.32), Inches(5.4), Inches(11), Inches(0.6),
      [("Las cifras salen de normativa oficial verificada. La IA solo redacta texto; nunca inventa un número.",
        15, RGBColor(0xC9,0xD6,0xE0), False, False, 0)])
_notes(s, "Cierre. Mensaje que se llevan: el sistema compara la calidad regulatoria del gas entre 21 jurisdicciones "
          "con trazabilidad total y sin cifras inventadas. Las cifras salen de normativa oficial verificada; la IA "
          "solo redacta texto.")

prs.save(PPTX)
print(f"PPTX generado: {os.path.relpath(PPTX, os.path.dirname(AQUI))}  ({os.path.getsize(PPTX)//1024} KB)  ·  {len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas")
