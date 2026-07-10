# -*- coding: utf-8 -*-
"""Genera las DOS presentaciones en PowerPoint EDITABLE (.pptx), estilo Enagás.

    python docs/generar_presentaciones_pptx.py

Salida:
    docs/Presentacion_7min.pptx   — versión ejecutiva (8 diapositivas ≈ 7 min)
    docs/Presentacion_15min.pptx  — versión completa (17 diapositivas ≈ 15 min)

Mismo contenido que los PDF homónimos, pero como PowerPoint nativo (cajas de texto,
tablas y formas editables) + NOTAS DEL PONENTE en cada diapositiva. Cubre todo el
proyecto: gas natural (21 jurisdicciones × 10 parámetros) + ampliación a biometano e
hidrógeno (dominio de red). Soporta **negrita** en línea con el marcador `**...**`.
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

AQUI = os.path.dirname(os.path.abspath(__file__))

# --- Paleta Enagás ---
AZUL   = RGBColor(0x01, 0x3A, 0x57)
CYAN   = RGBColor(0x00, 0x99, 0xD6)
VERDE  = RGBColor(0x6C, 0xB3, 0x3E)
GRIS   = RGBColor(0x4A, 0x5B, 0x68)
TINTA  = RGBColor(0x1B, 0x2A, 0x38)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRISCL = RGBColor(0xEE, 0xF2, 0xF6)
CLARO  = RGBColor(0xC9, 0xD6, 0xE0)
FUENTE = "Calibri"
PIE = "Comparador Regulatorio de Calidad de Gas · Gas natural · Biometano · Hidrógeno — Enagás"


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _box(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def _runs(p, text, size, color, bold=False):
    """Añade runs a un párrafo interpretando el marcador **negrita**."""
    for seg in re.split(r"(\*\*.*?\*\*)", text):
        if not seg:
            continue
        b = seg.startswith("**") and seg.endswith("**")
        r = p.add_run(); r.text = seg[2:-2] if b else seg
        r.font.name = FUENTE; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.bold = bool(b or bold)


def _text(slide, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """paras: lista de (texto, size, color, bold)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, r in enumerate(paras):
        txt, size, color, bold = (r + (False,))[:4]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        _runs(p, txt, size, color, bold)
    return tb


def _cabecera(slide, titulo, n):
    _box(slide, 0, 0, Inches(13.333), Inches(1.02), fill=AZUL)
    _box(slide, 0, Inches(1.02), Inches(13.333), Pt(4), fill=CYAN)
    _text(slide, Inches(0.55), Inches(0.14), Inches(12), Inches(0.78),
          [(titulo, 26, BLANCO, True)], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(0.55), Inches(7.08), Inches(11), Inches(0.32),
          [(PIE, 9.5, GRIS, False)])
    _text(slide, Inches(12.3), Inches(7.08), Inches(0.8), Inches(0.32),
          [(str(n), 9.5, GRIS, False)], align=PP_ALIGN.RIGHT)


def slide_bullets(prs, titulo, n, intro=None, bullets=None, notas="", bsize=17, space=10):
    s = _slide(prs); _cabecera(s, titulo, n)
    top = 1.4
    if intro:
        _text(s, Inches(0.6), Inches(top), Inches(12.1), Inches(0.9), [(intro, 15, GRIS, False)])
        top += 0.78
    tb = s.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(5.4))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in (bullets or []):
        sub = False
        if isinstance(b, tuple):
            b, sub = b
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space); p.level = 1 if sub else 0
        rb = p.add_run(); rb.text = "•  "
        rb.font.name = FUENTE; rb.font.size = Pt(bsize - 2 if sub else bsize)
        rb.font.color.rgb = GRIS if sub else CYAN
        _runs(p, b, (bsize - 2) if sub else bsize, GRIS if sub else TINTA)
    if notas:
        _notes(s, notas)
    return s


def slide_tabla(prs, titulo, n, headers, filas, notas="", intro=None, col_ratios=None,
                fsize=13, hsize=14, subtexto=None):
    s = _slide(prs); _cabecera(s, titulo, n)
    top = 1.45
    if intro:
        _text(s, Inches(0.6), Inches(top), Inches(12.1), Inches(0.8), [(intro, 15, GRIS, False)])
        top += 0.7
    nrows, ncols = len(filas) + 1, len(headers)
    tw = Inches(12.1)
    gr = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(top), tw, Inches(0.5)).table
    gr.first_row = False; gr.horz_banding = False
    if col_ratios:
        tot = sum(col_ratios)
        for j, cr in enumerate(col_ratios):
            gr.columns[j].width = int(tw * cr / tot)
    for j, htxt in enumerate(headers):
        c = gr.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = AZUL
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = Pt(3); c.margin_bottom = Pt(3)
        p = c.text_frame.paragraphs[0]; _runs(p, htxt, hsize, BLANCO, True)
    for i, fila in enumerate(filas, start=1):
        for j, val in enumerate(fila):
            c = gr.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = GRISCL if i % 2 else BLANCO
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = Pt(3); c.margin_bottom = Pt(3)
            p = c.text_frame.paragraphs[0]; p.word_wrap = True
            _runs(p, str(val), fsize, TINTA)
    if subtexto:
        _text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.8),
              [(subtexto, 11.5, GRIS, False)])
    if notas:
        _notes(s, notas)
    return s


def slide_portada(prs, sub, detalle, etiqueta, notas=""):
    s = _slide(prs)
    _box(s, 0, 0, Inches(13.333), Inches(7.5), fill=AZUL)
    _box(s, 0, Inches(4.55), Inches(13.333), Pt(5), fill=CYAN)
    _box(s, Inches(0.9), Inches(2.0), Inches(0.16), Inches(2.2), fill=VERDE)
    _text(s, Inches(1.25), Inches(2.0), Inches(11), Inches(2.3), [
        ("Comparador Regulatorio de", 40, BLANCO, True),
        ("Calidad de Gas", 40, BLANCO, True),
    ], space=2)
    _text(s, Inches(1.28), Inches(4.7), Inches(11.5), Inches(1.2), [
        (sub, 20, CYAN, False),
        (detalle, 15, CLARO, False),
    ], space=6)
    _text(s, Inches(1.28), Inches(6.5), Inches(11), Inches(0.5), [(etiqueta, 13, RGBColor(0x9F, 0xB3, 0xC4), False)])
    if notas:
        _notes(s, notas)
    return s


def slide_cierre(prs, lineas, pie, notas=""):
    s = _slide(prs)
    _box(s, 0, 0, Inches(13.333), Inches(7.5), fill=AZUL)
    _box(s, Inches(0.9), Inches(2.4), Inches(0.16), Inches(2.6), fill=VERDE)
    paras = [("En síntesis", 22, CYAN, True)] + [(ln, 26, BLANCO, True) for ln in lineas]
    _text(s, Inches(1.3), Inches(2.5), Inches(11), Inches(2.6), paras, space=8)
    _text(s, Inches(1.32), Inches(5.5), Inches(11), Inches(1.0), [(pie, 15, CLARO, False)])
    if notas:
        _notes(s, notas)
    return s


def slide_cifras(prs, n, notas=""):
    s = _slide(prs); _cabecera(s, "El sistema en cifras", n)
    tarjetas = [("21", "jurisdicciones"), ("10", "parámetros de calidad"),
                ("210", "valores verificables"), ("0", "cifras inventadas")]
    x = 0.6
    for num, txt in tarjetas:
        _box(s, Inches(x), Inches(2.2), Inches(2.95), Inches(2.5), fill=GRISCL)
        _box(s, Inches(x), Inches(2.2), Inches(2.95), Pt(6), fill=CYAN)
        _text(s, Inches(x), Inches(2.6), Inches(2.95), Inches(1.3), [(num, 54, AZUL, True)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(x), Inches(3.9), Inches(2.95), Inches(0.7), [(txt, 15, GRIS, False)], align=PP_ALIGN.CENTER)
        x += 3.08
    _text(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(0.6),
          [("Gas natural: 176 valores VERIFICADOS y 34 NO VERIFICABLES (la norma no los fija) — nunca inventados.", 15, TINTA, True)],
          align=PP_ALIGN.CENTER)
    _text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6),
          [("Ampliación a biometano e hidrógeno (dominio de red): España · Portugal · Francia · UE, con la misma disciplina.", 13, GRIS, False)],
          align=PP_ALIGN.CENTER)
    if notas:
        _notes(s, notas)
    return s


def slide_arquitectura(prs, n, notas=""):
    s = _slide(prs); _cabecera(s, "Arquitectura general", n)
    _text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
          [("Cuatro componentes con responsabilidades bien delimitadas.", 15, GRIS, False)])

    def comp(l, t, w, h, tit, sub, fill, tcol=BLANCO):
        _box(s, Inches(l), Inches(t), Inches(w), Inches(h), fill=fill)
        _text(s, Inches(l), Inches(t + 0.1), Inches(w), Inches(h - 0.15),
              [(tit, 16, tcol, True), (sub, 11.5, tcol, False)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=3)

    def flecha(l, t, w, h):
        a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
        a.fill.solid(); a.fill.fore_color.rgb = CYAN; a.line.fill.background(); a.shadow.inherit = False

    comp(4.9, 2.05, 3.5, 0.85, "INTERFAZ WEB", "El usuario formula su consulta", CYAN)
    flecha(6.45, 2.97, 0.4, 0.5)
    comp(4.4, 3.55, 4.5, 0.95, "SERVIDOR DE APLICACIÓN", "FastAPI · recibe, aplica la lógica y orquesta", AZUL)
    flecha(3.7, 4.6, 0.35, 0.5); flecha(9.25, 4.6, 0.35, 0.5)
    comp(1.35, 5.2, 4.3, 1.15, "BASE DE CONOCIMIENTO", "Ontología: valores verificados + sus fuentes", VERDE)
    comp(7.65, 5.2, 4.3, 1.15, "SERVICIO DE IA (externo)", "OpenAI: solo texto abierto · no genera cifras", GRIS)
    if notas:
        _notes(s, notas)
    return s


def slide_estados(prs, n, notas=""):
    s = _slide(prs); _cabecera(s, "Estados de verificación", n)
    _text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
          [("La garantía frente a la invención de datos: cada cifra está en uno de dos estados.", 15, GRIS, False)])
    _box(s, Inches(0.7), Inches(2.15), Inches(5.8), Inches(2.7), fill=GRISCL)
    _box(s, Inches(0.7), Inches(2.15), Inches(5.8), Pt(6), fill=VERDE)
    _text(s, Inches(1.0), Inches(2.4), Inches(5.2), Inches(2.4), [
        ("VERIFICADO", 22, AZUL, True),
        ("176 valores (gas natural)", 15, VERDE, True),
        ("Cifra contrastada literalmente (verbatim) contra su boletín oficial.", 14, TINTA, False),
    ], space=8)
    _box(s, Inches(6.85), Inches(2.15), Inches(5.8), Inches(2.7), fill=GRISCL)
    _box(s, Inches(6.85), Inches(2.15), Inches(5.8), Pt(6), fill=CYAN)
    _text(s, Inches(7.15), Inches(2.4), Inches(5.2), Inches(2.4), [
        ("NO VERIFICABLE", 22, AZUL, True),
        ("34 valores", 15, CYAN, True),
        ("La norma de esa jurisdicción no fija ese parámetro. No se estima: se marca y se explica.", 14, TINTA, False),
    ], space=8)
    _text(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.0),
          [("No hay estado intermedio: el valor consta en la norma, o se declara que la norma no lo establece.", 15, TINTA, True)],
          align=PP_ALIGN.CENTER)
    if notas:
        _notes(s, notas)
    return s


# ================================ 7 MINUTOS ==================================
def build_7min(path):
    prs = _new_prs()
    slide_portada(prs,
                  "Gas natural · Biometano · Hidrógeno",
                  "Comparativa regulatoria trazable entre jurisdicciones — sin cifras inventadas",
                  "Presentación ejecutiva · 7 minutos",
                  notas="Objetivo: explicar en 7 minutos qué hace la herramienta, cómo está construida y por qué es fiable. "
                        "Mensaje que se llevan desde el primer minuto: compara con trazabilidad total y sin inventar cifras.")
    slide_bullets(prs, "El problema", 2,
                  intro="Comparar la calidad regulatoria del gas entre países es hoy manual y arriesgado.",
                  bullets=[
                      "Cada país regula la calidad admisible con **su propia normativa**: poder calorífico, índice de Wobbe, azufre, CO₂, puntos de rocío…",
                      "Están **dispersas** en boletines oficiales distintos, en varios idiomas y con **unidades y condiciones de referencia diferentes**.",
                      "Compararlas a mano es **laborioso y propenso a error**.",
                      "Y ahora se suman **biometano e hidrógeno**, con marcos regulatorios aún en construcción.",
                  ],
                  notas="Cada país regula la calidad admisible del gas con su propia normativa. Está dispersa, en varios idiomas y "
                        "con unidades distintas; compararla a mano es costoso y arriesgado. Y ahora se suman biometano e hidrógeno. "
                        "Ese es el hueco que cubre la herramienta.")
    slide_cifras(prs, 3,
                 notas="Cuatro cifras resumen el alcance en gas natural: 21 jurisdicciones, 10 parámetros, 210 valores y cero "
                       "inventados. De los 210, 176 verificados verbatim y 34 marcados como no verificable. La ampliación a "
                       "biometano e hidrógeno usa la misma disciplina para España, Portugal, Francia y la UE.")
    slide_arquitectura(prs, 4,
                       notas="Cuatro componentes: la interfaz web, el servidor de aplicación (FastAPI, el cerebro), la base de "
                             "conocimiento (ontología con los datos verificados) y un servicio de IA externo, acotado, que solo "
                             "redacta texto y nunca genera cifras. Los datos y cálculos son nuestros; la IA es un auxiliar controlado.")
    slide_bullets(prs, "El principio: cero cifras inventadas", 5,
                  intro="El sistema combina dos mundos separados a propósito.",
                  bullets=[
                      "**Mundo determinista** (código + ontología): la **única** fuente de cifras, límites, conversiones y comparaciones. Nunca improvisa un número.",
                      "**Mundo conversacional** (IA): interpreta la pregunta y **redacta**, pero tiene **prohibido generar cifras**; las obtiene llamando a herramientas deterministas.",
                      "Regla: **determinista-primero, IA-como-respaldo**. Lo cuantitativo lo resuelve el código; solo el texto abierto pasa a la IA.",
                      ("Cada valor cita norma, artículo, página y enlace: trazabilidad completa.", True),
                  ],
                  notas="El corazón del proyecto. Conviven dos mundos: el determinista (código + ontología) es la única fuente de "
                        "cifras y nunca improvisa; el conversacional (IA) redacta pero tiene prohibido generar cifras. Regla: "
                        "determinista primero, IA como respaldo. Cada valor sale citando su fuente.")
    slide_bullets(prs, "Qué puede hacer el sistema", 6,
                  intro="Cinco secciones, todas con la misma garantía de cero cifras inventadas:",
                  bullets=[
                      "**Consulta libre** (chat), con **análisis de interconexión en cadena**: para una ruta (p. ej. España-Francia-Alemania) halla el **cuello de botella** regulatorio y avisa de incompatibilidades.",
                      "**Comparativa**: parámetro a parámetro entre países y **matriz** (heatmap) países × parámetros, con **exportación a Excel/PDF**.",
                      "**Analizar gas**: composición de un gas y respuesta país a país (**cumple / alerta / no cumple**), con la cita de cada límite.",
                      "**Comparativa de biometano** y de **hidrógeno**: el mismo motor, para los dos gases nuevos.",
                  ], bsize=15,
                  notas="Cinco secciones con la misma garantía. Consulta libre (chat) con análisis de interconexión en cadena, que "
                        "identifica el cuello de botella regulatorio de una ruta. Comparativa con matriz y exportación. Analizar gas, "
                        "que valida un gas país a país. Y las dos secciones nuevas de biometano e hidrógeno, con el mismo motor.")
    slide_tabla(prs, "Ampliación: biometano e hidrógeno", 7,
                headers=["Gas", "Jurisdicciones", "Parámetros clave", "Fuentes normativas"],
                filas=[
                    ["**Biometano**", "España · Portugal · Francia · UE", "CH₄ mínimo · CO₂ · siloxanos (+ impurezas)",
                     "EN 16723-1 · EN 16726 · Reg. (UE) 2024/1789 · Orden TED/181/2025 · GRTgaz · RQS (PT)"],
                    ["**Hidrógeno**", "Dominio de RED (TSO) + producto", "Pureza H₂ · O₂ · trazas de compresores",
                     "CEN/TS 17977 · GIE (recom.) · RQS Anexo XII (PT, 98 % vinculante) · ISO 14687 (vehículo)"],
                ],
                intro="Se añade como capa aditiva: el gas natural queda intacto (mismo motor, misma disciplina de verificación).",
                col_ratios=[1.3, 2.4, 2.6, 4.2], fsize=11,
                subtexto="Distinción clave: calidad de RED (gasoducto: CEN/TS 17977, GIE, ≥98 %) frente a producto de VEHÍCULO (ISO 14687, 99,97 %). "
                         "Un estudio de terminología justifica una capa de búsqueda semántica, preparada y opt-in.",
                notas="La ampliación es capa aditiva: el gas natural queda intacto. Biometano cubre España, Portugal, Francia y la UE. "
                      "En hidrógeno, la distinción clave: calidad de RED (gasoducto, CEN/TS 17977, GIE) frente a producto de VEHÍCULO "
                      "(ISO 14687). Solo Portugal fija hoy una pureza vinculante del 98 %.")
    slide_cierre(prs,
                 ["Compara la calidad regulatoria del gas natural,",
                  "biometano e hidrógeno con trazabilidad total",
                  "y sin cifras inventadas."],
                 "Las cifras salen de normativa oficial verificada. La IA solo redacta texto; nunca inventa un número.",
                 notas="Cierre. Mensaje: compara la calidad regulatoria de los tres gases con trazabilidad total y sin cifras inventadas.")
    prs.save(path)
    return path, len(prs.slides._sldIdLst)


# ================================ 15 MINUTOS =================================
def build_15min(path):
    prs = _new_prs()
    slide_portada(prs,
                  "Gas natural · Biometano · Hidrógeno",
                  "Arquitectura, funcionamiento y garantías — comparativa entre jurisdicciones",
                  "Presentación completa · 15 minutos",
                  notas="Presentación de 15 minutos: arquitectura, funcionamiento y garantías. Hilo conductor: trazabilidad total y "
                        "cero cifras inventadas. Nació para gas natural y se ha extendido a biometano e hidrógeno.")
    slide_bullets(prs, "Índice", 2,
                  intro="Recorrido de la presentación:",
                  bullets=[
                      "Contexto y objetivo · el sistema en cifras",
                      "Arquitectura general y funcionalidades",
                      "Las tres capas de datos y la ontología (base de conocimiento)",
                      "Estados de verificación · FastAPI frente a la API de OpenAI",
                      "El motor determinista y la normalización (ISO 13443)",
                      "La inteligencia artificial (con salvaguardas) y la recuperación documental (RAG)",
                      "Ampliación a biometano e hidrógeno · estudio de terminología",
                      "Metodología, verificación y garantías",
                  ], bsize=16, space=8,
                  notas="Agenda. El hilo es por componentes: primero visión general y cifras, luego cada pieza (datos, motor, IA), y "
                        "por último la ampliación, la metodología y las garantías.")
    slide_bullets(prs, "Contexto y objetivo", 3,
                  intro="El problema y la propuesta de valor.",
                  bullets=[
                      "Cada país regula la calidad admisible del gas con **su propia normativa** (poder calorífico, Wobbe, azufre, CO₂, puntos de rocío…).",
                      "Está **dispersa**, en varios idiomas y con **unidades y condiciones de referencia distintas**; compararla a mano es costoso y arriesgado.",
                      "**Solución:** un asistente que compara esa calidad entre **21 jurisdicciones y 10 parámetros**, en lenguaje natural.",
                      "**Principio de diseño:** ausencia de cifras no verificadas; todo valor procede de normativa oficial, con su cita.",
                      "Ampliado a **biometano e hidrógeno**, manteniendo el gas natural intacto.",
                  ],
                  notas="El problema: normativa dispersa, en varios idiomas y con unidades distintas. La solución: un asistente que "
                        "compara 21 jurisdicciones y 10 parámetros en lenguaje natural, con el principio de no generar cifras por "
                        "estimación. Ampliado a biometano e hidrógeno sin tocar el gas natural.")
    slide_cifras(prs, 4,
                 notas="Cuatro cifras fijan el alcance. Detrás del cero está lo importante: de 210 valores, 176 verificados verbatim "
                       "y 34 declarados no verificable porque la norma no los fija. No se rellenan con estimaciones.")
    slide_arquitectura(prs, 5,
                       notas="Cuatro componentes. La interfaz web, el servidor FastAPI que concentra la lógica, la ontología con los "
                             "datos verificados y el servicio de IA externo, acotado. Los datos y cálculos son propios; la IA es un "
                             "proveedor auxiliar bajo control.")
    slide_bullets(prs, "Funcionalidades", 6,
                  intro="Cinco secciones, todas con la misma garantía de cero cifras inventadas:",
                  bullets=[
                      "**Consulta libre** (chat): límites, cumplimiento, fuentes y comparaciones en lenguaje natural. Incluye **interconexión en cadena**: para una ruta de varios países halla el **cuello de botella** regulatorio.",
                      "**Comparativa**: comparación puntual de un parámetro entre países y **matriz** (heatmap) países × parámetros, con **exportación a Excel/PDF**.",
                      "**Analizar gas**: composición de un gas concreto → respuesta país a país (**cumple / alerta / no cumple**) con la cita de cada límite.",
                      "**Comparativa de biometano** y de **hidrógeno**: el mismo motor y la misma presentación, para los gases nuevos.",
                  ], bsize=15, space=9,
                  notas="Cinco funcionalidades. La consulta libre incluye el análisis de interconexión en cadena. La comparativa da la "
                        "matriz completa con exportación. Analizar gas valida un gas país a país. Y las secciones de biometano e "
                        "hidrógeno replican la experiencia con el mismo motor.")
    slide_tabla(prs, "Las tres capas de datos", 7,
                headers=["Capa", "Contenido", "Función"],
                filas=[
                    ["**1. Documentos oficiales**", "Los PDF de las normas (BOE, ERSE, DVGW, Fluxys, National Grid, RQS…)", "Fuente primaria y última de verdad"],
                    ["**2. Ontología**", "Fichero estructurado con las cifras extraídas de esos PDF, con su contexto", "Repositorio del que salen las respuestas"],
                    ["**3. Índice documental (RAG)**", "Índice del texto de los PDF, segmentado en fragmentos con solape", "Buscador interno para consultas abiertas"],
                ],
                intro="No hay una única base de datos gigante: la información se organiza en tres capas con funciones distintas.",
                col_ratios=[2.2, 4.3, 3.1],
                notas="Las cifras residen en la capa 2, la ontología, y cada una referencia su documento oficial de la capa 1. La capa "
                      "3 no almacena ninguna cifra: es un índice de búsqueda sobre el texto, para las consultas abiertas.")
    slide_bullets(prs, "La base de conocimiento (ontología)", 8,
                  intro="El elemento central. De cada valor se registra su cifra y todo su contexto normativo:",
                  bullets=[
                      "El **valor** (o su rango) y la **unidad**.",
                      "Las **condiciones de referencia** (a qué temperatura y presión se mide).",
                      "El **texto literal** de la norma, tal como está redactado.",
                      "La **cita completa**: norma, artículo, página y enlace.",
                      "Una **nota** aclaratoria con los matices, y el **estado de verificación**.",
                      ("Se usa un fichero YAML (legible y auditable) en lugar de una base de datos relacional: a esta escala es más trazable y se versiona junto al código.", True),
                  ],
                  notas="La ontología guarda de cada dato no solo el número, sino todo su contexto: unidad, condiciones, texto literal, "
                        "cita completa, nota y estado de verificación. Se usa YAML, legible y auditable, versionado junto al código.")
    slide_estados(prs, 9,
                  notas="La garantía anti-invención. Verificado: 176 valores contrastados verbatim. No verificable: 34 que la norma no "
                        "fija; no se completan con una estimación. Ejemplo: en Dinamarca los límites de O₂/CO₂ son de biogás, no de gas "
                        "natural, así que para gas natural se dejaron como no verificable.")
    slide_tabla(prs, "FastAPI y la API de OpenAI: son cosas distintas", 10,
                headers=["", "FastAPI", "API de OpenAI"],
                filas=[
                    ["**Naturaleza**", "Framework para construir NUESTRO servidor", "Servicio externo que consumimos"],
                    ["**Titularidad**", "Propia (es nuestro backend)", "De OpenAI (somos cliente)"],
                    ["**Coste**", "Sin coste (código abierto)", "De pago, por uso"],
                    ["**Papel**", "Núcleo de la aplicación", "Proveedor auxiliar, invocado de forma controlada"],
                ],
                intro="Ambos incluyen el término «API», pero son componentes de niveles distintos.",
                col_ratios=[1.5, 4.0, 4.0],
                notas="FastAPI es el framework con el que hemos hecho nuestro servidor: infraestructura propia y gratuita. La API de "
                      "OpenAI es un servicio de terceros que consumimos puntualmente. Nuestro servidor es el imprescindible: atiende la "
                      "web, accede a los datos, calcula, decide cuándo usar la IA y custodia las credenciales.")
    slide_bullets(prs, "El motor determinista", 11,
                  intro="Toda consulta pasa primero por un enrutado. «Determinista» = misma consulta, misma respuesta, por código, sin IA.",
                  bullets=[
                      "Consultas **cuantitativas** (un límite, un cumplimiento, una comparación, una conversión): las resuelve el **código leyendo la ontología**. Sin IA, sin posibilidad de generar un valor incorrecto.",
                      "Consultas de **texto abierto** («¿en qué consiste el índice de Wobbe?»): se derivan al servicio de IA.",
                      "Resuelve **sin IA** siete tipos de intención: valor de un límite, cumplimiento, fuente, intercambiabilidad, restrictividad frente a España, comparación directa y conversión de condiciones.",
                      ("En la práctica, la mayoría de consultas se resuelven sin recurrir a la IA.", True),
                  ],
                  notas="Toda consulta pasa por el enrutado. Las cuantitativas las resuelve el código leyendo la ontología, sin IA. Las "
                        "de texto abierto se derivan a la IA. El motor resuelve por sí solo siete tipos de intención; la mayoría de "
                        "consultas no llegan a la IA.")
    slide_bullets(prs, "Normalización de condiciones (ISO 13443)", 12,
                  intro="Para comparar de forma rigurosa: cada país usa unidades y condiciones de referencia distintas.",
                  bullets=[
                      "Unos expresan en **kWh/m³**, otros en **MJ/m³** o kcal/m³; unos miden a **0 °C**, otros a **15** o a **25 °C**.",
                      "Comparar los valores en bruto sería incorrecto (como comparar millas con kilómetros).",
                      "Todos los valores se llevan a la **base española** con los **factores literales de la norma ISO 13443** (Tabla A.1), ya implementados y verificados.",
                      "**España** es siempre la base de referencia de la comparación.",
                      ("Los valores derivados se muestran con 2 decimales (sin falsa precisión); los originales, con su precisión de origen.", True),
                  ],
                  notas="Cada país expresa sus límites en unidades y condiciones distintas; compararlos en bruto sería incorrecto. Todos "
                        "los valores se llevan a la base española con los factores de la ISO 13443, que se toman literalmente y están "
                        "verificados. España es siempre la referencia.")
    slide_bullets(prs, "La inteligencia artificial y sus salvaguardas", 13,
                  intro="Cuando una consulta se deriva al servicio de IA, este opera bajo restricciones estrictas:",
                  bullets=[
                      "Tiene **prohibido generar cifras**: si necesita un dato, lo pide a las herramientas internas, que lo obtienen de la ontología.",
                      "Debe **citar** los documentos oficiales, y su ámbito se limita a la calidad del gas.",
                      "Si el servicio de IA no está disponible, el sistema **conmuta automáticamente al modo determinista**: nunca se interrumpe.",
                      ("Modelo: OpenAI GPT-4o-mini, temperatura 0 (máxima previsibilidad), con llamadas a herramientas (function calling).", True),
                  ],
                  notas="Cuando la consulta llega a la IA, opera atada en corto: prohibido generar cifras (las pide a las herramientas), "
                        "debe citar y se limita a la calidad del gas. Si la IA no está disponible, conmuta al modo determinista y nunca "
                        "se interrumpe.")
    slide_bullets(prs, "Recuperación documental (RAG) y terminología", 14,
                  intro="Para el texto abierto, la respuesta se fundamenta en los documentos oficiales, no en el conocimiento general del modelo.",
                  bullets=[
                      "**Indexación:** se procesan los PDF, se extrae el texto y se trocea en **fragmentos con solape** (ventana continua que cruza el salto de página); indexación **incremental**.",
                      "**Recuperación:** búsqueda **léxica** (por términos) que devuelve los fragmentos pertinentes con archivo y página; reproducible, sin caja negra externa.",
                      "**Estudio de terminología:** se mide la variación de nombres entre normas (**índice de variación**: 27,4 en gas natural, 9,2 en biometano, 7,3 en hidrógeno; umbral 7,0).",
                      ("El estudio justifica una capa de búsqueda semántica (multilingüe), que se deja preparada y opt-in.", True),
                  ], bsize=15,
                  notas="El RAG fundamenta las respuestas de texto abierto en los documentos oficiales. Indexación en fragmentos con "
                        "solape; recuperación léxica reproducible. Además, un estudio de terminología mide la variación de nombres entre "
                        "normas y justifica una capa de búsqueda semántica, que se deja preparada y activable a voluntad.")
    slide_tabla(prs, "Ampliación: biometano e hidrógeno", 15,
                headers=["Gas", "Jurisdicciones", "Parámetros clave", "Fuentes normativas"],
                filas=[
                    ["**Biometano**", "España · Portugal · Francia · UE", "CH₄ mínimo · CO₂ · siloxanos (+ O₂, azufre, NH₃, aminas)",
                     "EN 16723-1 · EN 16726 · Reg. (UE) 2024/1789 · Orden TED/181/2025 · GRTgaz/GRDF · RQS Anexo XI"],
                    ["**Hidrógeno**", "Dominio de RED (TSO) + producto", "Pureza H₂ · O₂ · trazas de compresores",
                     "CEN/TS 17977 · GIE (recomendación) · RQS Anexo XII (PT, 98 % vinculante) · ISO 14687 (vehículo)"],
                ],
                intro="Capa aditiva: el gas natural queda byte a byte intacto (mismo motor, misma disciplina). 14/14 pruebas en verde.",
                col_ratios=[1.2, 2.3, 2.9, 4.3], fsize=10.5,
                subtexto="Distinción de dominio: calidad de RED (gasoducto: CEN/TS 17977, GIE, ≥98 %) frente a producto de VEHÍCULO (ISO 14687, 99,97 %). "
                         "Solo Portugal (RQS) fija hoy una pureza de H₂ vinculante; España y Francia regulan el blending; la UE lo recomienda vía GIE.",
                notas="La ampliación es capa aditiva: el gas natural queda intacto, lo confirman las pruebas. Distinción de dominio: "
                      "calidad de RED (gasoducto, CEN/TS 17977, GIE) frente a producto de VEHÍCULO (ISO 14687). Hoy solo Portugal fija "
                      "una pureza vinculante del 98 %.")
    slide_bullets(prs, "Metodología, verificación y garantías", 16,
                  intro="La fiabilidad es consecuencia del rigor del proceso de carga. Para cada jurisdicción:",
                  bullets=[
                      "Se identificó la **norma oficial vigente**, se **archivó el documento** localmente y se **transcribió cada cifra literalmente**.",
                      "Se **verificó una a una** contra el documento; lo que la norma **no fija, no se completa** (se marca «no verificable», con su justificación).",
                      "**Controles automáticos**: comprueban que las celdas se resuelven, que los enlaces funcionan y que no hay incoherencias.",
                      ("Garantías: cero cifras inventadas · trazabilidad completa · transparencia · reproducibilidad · auditabilidad.", True),
                  ],
                  notas="La fiabilidad es consecuencia del rigor: identificar la norma vigente, archivar el documento, transcribir cada "
                        "cifra literalmente, verificarla una a una y marcar como no verificable lo que la norma no fija. Con controles "
                        "automáticos. De ahí las cinco garantías.")
    slide_cierre(prs,
                 ["Compara la calidad regulatoria del gas natural,",
                  "biometano e hidrógeno entre jurisdicciones,",
                  "con trazabilidad total y sin cifras inventadas."],
                 "Las cifras salen de normativa oficial verificada. La IA solo redacta texto; nunca inventa un número. El gas natural permanece intacto tras la ampliación.",
                 notas="Cierre. La herramienta compara la calidad regulatoria de los tres gases con trazabilidad total y sin cifras "
                       "inventadas, y la ampliación no ha degradado nada de lo anterior.")
    prs.save(path)
    return path, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    for fn, name in [(build_7min, "Presentacion_7min.pptx"), (build_15min, "Presentacion_15min.pptx")]:
        p, ns = fn(os.path.join(AQUI, name))
        print(f"PPTX generado: docs/{name}  ({os.path.getsize(p)//1024} KB · {ns} diapositivas)")
